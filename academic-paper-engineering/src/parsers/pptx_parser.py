"""PPTX 文档解析器 - 将 PPT 幻灯片解析为资产 IR"""

from pathlib import Path
from typing import Dict, List


class PptxParser:
    """解析 .pptx 格式的演示文稿，提取可用的论文资产"""

    def __init__(self):
        self.ir = {
            "source_file": "",
            "slides": []
        }

    def parse(self, file_path: str) -> Dict:
        """解析 PPTX 文件，返回资产 IR"""
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        self.ir["source_file"] = path.name

        try:
            from pptx import Presentation
            prs = Presentation(str(path))

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_data = {
                    "slide_number": slide_num,
                    "title": "",
                    "assets": [],
                    "text_content": ""
                }

                texts = []
                for shape in slide.shapes:
                    # 提取标题
                    if shape.has_text_frame:
                        if shape == slide.shapes.title:
                            slide_data["title"] = shape.text_frame.text.strip()
                        texts.append(shape.text_frame.text.strip())

                    # 提取图片
                    if shape.shape_type == 13:  # PICTURE
                        slide_data["assets"].append({
                            "type": "figure",
                            "source": f"slide{slide_num}_image.png",
                            "description": "",
                            "suggested_caption": ""
                        })

                    # 提取表格
                    if shape.has_table:
                        table = shape.table
                        headers = []
                        rows = []
                        for row_idx, row in enumerate(table.rows):
                            cells = [cell.text.strip() for cell in row.cells]
                            if row_idx == 0:
                                headers = cells
                            else:
                                rows.append(cells)
                        slide_data["assets"].append({
                            "type": "table",
                            "data": {
                                "headers": headers,
                                "rows": rows
                            },
                            "suggested_caption": ""
                        })

                slide_data["text_content"] = '\n'.join(texts)
                self.ir["slides"].append(slide_data)

        except ImportError:
            raise ImportError("需要安装 python-pptx: pip install python-pptx")

        return self.ir

    def extract_images(self, file_path: str, output_dir: str) -> List[str]:
        """提取 PPTX 中的所有图片到指定目录"""
        path = Path(file_path)
        output = Path(output_dir)
        output.mkdir(parents=True, exist_ok=True)

        extracted = []
        try:
            from pptx import Presentation
            prs = Presentation(str(path))

            img_counter = 0
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if shape.shape_type == 13:
                        img_counter += 1
                        image = shape.image
                        ext = image.ext
                        img_path = output / f"slide{slide_num}_image{img_counter}.{ext}"
                        with open(img_path, 'wb') as f:
                            f.write(image.blob)
                        extracted.append(str(img_path))

        except ImportError:
            raise ImportError("需要安装 python-pptx: pip install python-pptx")

        return extracted
